from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsmap
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Colors matching website
GREEN = RGBColor(74, 156, 109)
DARK = RGBColor(26, 26, 26)
GRAY = RGBColor(100, 100, 100)
LIGHT_GREEN_BG = RGBColor(241, 248, 233)
LIGHT_PINK_BG = RGBColor(252, 228, 236)

def add_title_slide(title, subtitle="", notes="", bg_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg_color:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.name = "Georgia"
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.name = "Arial"
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_content_slide(label, title, bullets=None, text=None, images=None, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = Inches(0.6)
    
    # Section label
    if label:
        lbl = slide.shapes.add_textbox(Inches(0.8), y, Inches(14), Inches(0.4))
        p = lbl.text_frame.paragraphs[0]
        p.text = label.upper()
        p.font.size = Pt(11)
        p.font.name = "Arial"
        p.font.color.rgb = GRAY
        p.font.bold = True
        y += Inches(0.35)
    
    # Title
    ttl = slide.shapes.add_textbox(Inches(0.8), y, Inches(14), Inches(1))
    p = ttl.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Georgia"
    p.font.color.rgb = DARK
    y += Inches(0.8)
    
    # Body text
    if text:
        tb = slide.shapes.add_textbox(Inches(0.8), y, Inches(9), Inches(1.2))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.name = "Arial"
        p.font.color.rgb = GRAY
        y += Inches(0.9)
    
    # Bullets
    if bullets:
        for b in bullets:
            tb = slide.shapes.add_textbox(Inches(0.8), y, Inches(9), Inches(0.6))
            tb.text_frame.word_wrap = True
            p = tb.text_frame.paragraphs[0]
            p.text = "→  " + b
            p.font.size = Pt(16)
            p.font.name = "Arial"
            p.font.color.rgb = GRAY
            y += Inches(0.4)
    
    # Images
    if images:
        img_x = Inches(10) if (text or bullets) else Inches(0.8)
        img_y = Inches(2) if (text or bullets) else y
        img_width = Inches(5) if (text or bullets) else Inches(14)
        
        if isinstance(images, str):
            images = [images]
        
        if len(images) == 1 and not (text or bullets):
            # Single image, centered
            if os.path.exists(images[0]):
                slide.shapes.add_picture(images[0], Inches(3), Inches(2.5), width=Inches(10))
        elif len(images) == 1:
            # Single image on right
            if os.path.exists(images[0]):
                slide.shapes.add_picture(images[0], Inches(10), Inches(2), width=Inches(5))
        elif len(images) == 2:
            # Two images side by side
            for i, img in enumerate(images):
                if os.path.exists(img):
                    slide.shapes.add_picture(img, Inches(0.8 + i*7.5), Inches(4.5), width=Inches(6.5))
        elif len(images) == 3:
            # Three images in a row
            for i, img in enumerate(images):
                if os.path.exists(img):
                    slide.shapes.add_picture(img, Inches(0.8 + i*5), Inches(4.5), width=Inches(4.5))
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_image_slide(label, title, images, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    if label:
        lbl = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(14), Inches(0.4))
        p = lbl.text_frame.paragraphs[0]
        p.text = label.upper()
        p.font.size = Pt(11)
        p.font.name = "Arial"
        p.font.color.rgb = GRAY
        p.font.bold = True
    
    if title:
        ttl = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(14), Inches(0.8))
        p = ttl.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.name = "Georgia"
        p.font.color.rgb = DARK
    
    if isinstance(images, str):
        images = [images]
    
    if len(images) == 1:
        if os.path.exists(images[0]):
            slide.shapes.add_picture(images[0], Inches(1.5), Inches(2), width=Inches(13))
    elif len(images) == 2:
        for i, img in enumerate(images):
            if os.path.exists(img):
                slide.shapes.add_picture(img, Inches(0.8 + i*7.5), Inches(2), width=Inches(7))
    elif len(images) == 3:
        for i, img in enumerate(images):
            if os.path.exists(img):
                slide.shapes.add_picture(img, Inches(0.5 + i*5.2), Inches(2), width=Inches(4.8))
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

# ============ SLIDES ============

# TITLE
add_title_slide("Portfolio Review", "Elise Nguyen\nUX Designer at Amazon Kindle",
    notes="Introduce yourself. Thank them for their time.")

# PROJECT 1: SHOP GOODREADS
add_title_slide("Shop Goodreads on Amazon", "Helping customers buy their next read", 
    bg_color=LIGHT_GREEN_BG,
    notes="Transition to first project. This was a 4-month project that launched in April 2025.")

add_content_slide("Overview", "The Project",
    text="I designed a highly-requested feature that brings users' Goodreads Want to Read shelf into Amazon Your Books, letting them browse, filter, and shop their WTR titles all in one place.",
    bullets=["Role: Primary UX Designer", "Timeline: 4 months; Launched April 2025", "Team: Kindle and Goodreads product and engineering teams"],
    images="assets/Project 1/thumbnail.png",
    notes="This was the first-ever collaboration between Goodreads and Amazon Your Books teams. I was the sole designer on the project.")

add_content_slide("Problem", "A disconnected shopping experience",
    text="Readers use Goodreads to track books they want to read, but this information isn't easily accessible when they shop on Amazon.",
    bullets=["Users struggle to find books already saved to their Want to Read shelf", "41.7M linked Goodreads and Amazon customers", "Opportunity to integrate Goodreads into Amazon for seamless shopping"],
    notes="The key insight here is that there's a huge user base already linking accounts, but the experience between the two platforms was fragmented.")

add_image_slide("", "", "assets/Project 1/image 1.png",
    notes="This visual shows the disconnect - users had to manually search for books they'd already saved on Goodreads.")

add_content_slide("Design Challenge", "How might we reduce friction between tracking books on Goodreads and purchasing them on Amazon?",
    notes="This was our north star question throughout the project.")

add_content_slide("The Challenge", "Balancing integration with trust",
    text="Data showed the separation between Goodreads and Amazon was intentional:",
    bullets=["Only one-third of linked users maintain lists on both platforms", "Just one-fifth of books overlap between platforms", "Users see Goodreads as impartial discovery, Amazon as transactional"],
    notes="This was a critical insight from research. Users deliberately kept these platforms separate because they trusted Goodreads recommendations to be unbiased. We had to be careful not to break that trust.")

add_content_slide("Mission", "Design Principles",
    bullets=["DO: Convenience - Single source of truth for total reading life", "DO: Trust - Control visibility of books on both platforms", "DO: Transparency - Context for where/why they see Goodreads content", "DON'T: Friction - No context-switching between apps", "DON'T: Bias - No sponsored content in Goodreads recommendations"],
    notes="These principles guided every design decision. Whenever we had debates, we came back to these.")

add_content_slide("Design Decision #1", "Information Architecture",
    text="Where does Goodreads fit into Your Books?",
    bullets=["Option 1: Add Goodreads WTR as a filter and card", "Option 2: Add Goodreads as a secondary tab", "Option 3: Add Goodreads as a primary tab"],
    images="assets/Project 1/image 22.png",
    notes="We explored three main approaches. Each had tradeoffs around discoverability, complexity, and future scalability.")

add_content_slide("", "Decision: Full integration",
    text="We aligned on fully integrating Goodreads Want to Read into the existing Your Books experience rather than introducing a separate Goodreads tab.",
    bullets=["Builds on customers' current mental models", "Avoids fragmenting the experience", "Enables WTR books to benefit from Your Books' filters and recommendations"],
    notes="This was a key alignment moment with both Amazon and Goodreads design teams. The decision prioritized user mental models over organizational boundaries.")

add_content_slide("Design Decision #2", "Badging",
    text="How do we differentiate Amazon wishlists from Goodreads shelves?",
    images="assets/Project 1/image 23.png",
    notes="We explored many options - logos, labels, colors. Each had tradeoffs around clarity vs. visual noise.")

add_content_slide("User Research", "Validating with users",
    text="Ran unmoderated usability test with 8 users on 3 badging options:",
    images="assets/Project 1/image 26.png",
    notes="User research was critical here because internal opinions were divided. We needed real user input to break the tie.")

add_content_slide("", "Research findings",
    text="Clear preference for 'On Goodreads shelf' label due to clarity and recognizability.",
    bullets=["Logo-only approaches felt cleaner but introduced discoverability risks", "Non-Goodreads users struggled with logo-only options", "Explicit labeling prioritized comprehension over visual minimalism"],
    images="assets/Project 1/goodreads-badging-final.png",
    notes="The key learning: when introducing something new, clarity beats elegance. We can evolve to more compact designs as users become familiar.")

add_image_slide("Solution", "Introducing a new integration",
    ["assets/Project 1/image 2.png", "assets/Project 1/image 3.png"],
    notes="A new modal within Your Books and email to Goodreads members introduces this integration, guiding them to the new experience.")

add_image_slide("", "See all your lists in one place",
    ["assets/Project 1/image 4.png", "assets/Project 1/image 5.png"],
    notes="In the Lists tab, users can now view their Goodreads Want to Read shelf alongside Amazon wishlists. A new Goodreads filter and branded badge help users easily identify and navigate.")

add_image_slide("", "Shop all your books at once",
    ["assets/Project 1/image 6.png", "assets/Project 1/image 7.png", "assets/Project 1/image 8.png"],
    notes="In the All Titles tab, users can browse their complete book collection from both Amazon and Goodreads in one unified view.")

add_image_slide("", "Your books. Your choice.",
    ["assets/Project 1/image 9.png", "assets/Project 1/image 10.png", "assets/Project 1/image 11.png"],
    notes="Users have full control to learn more about the feature, hide Goodreads books from their view, or unlink accounts entirely. This was critical for maintaining trust.")

add_content_slide("Outcome", "Impact",
    text="Launched April 22, 2025 for 9.2M linked monthly active users (41.7M total).",
    bullets=["86% higher engagement for WTR shelved books vs. Wishlist books", "Strong positive customer feedback", "Featured in Fast Company and viral TikTok coverage"],
    notes="The engagement numbers exceeded our expectations. The 86% higher engagement showed users were actively using the integration.")

add_content_slide("", "What users are saying",
    bullets=['"This has been something I\'ve wanted for years."', '"I don\'t have to go back and forth to remember what I put on my TBR list!"', '"No longer having to track my want to buy book list in both Amazon and Goodreads"'],
    images=["assets/Project 1/image 13.png", "assets/Project 1/tiktok-preview.png"],
    notes="These customer quotes validated our design decisions. The pain point we identified was real and the solution resonated.")

add_content_slide("Learnings", "What I learned",
    bullets=["Design across two distinct platforms and brands while navigating complex cross-team dynamics", "Drive alignment on decisions and synthesize feedback from many stakeholders", "Design within existing information architecture that differed from both Goodreads and Kindle", "Thoughtfully merge two brands by respecting what makes each valuable to users"],
    notes="This project taught me a lot about cross-team collaboration and designing for trust. The biggest challenge wasn't the UI - it was aligning two different organizations with different goals.")

# PROJECT 2: COLOR WRITING
add_title_slide("Color Writing on Kindle", "Designing color writing and annotation for Kindle Scribe's debut e-ink color display",
    bg_color=LIGHT_PINK_BG,
    notes="Transition to second project. This was a 7-month project for a hardware launch - very different pace and constraints than the first project.")

add_content_slide("Overview", "The Project",
    text="I owned the design for the color writing and annotation experience for Kindle Scribe's debut e-ink color display.",
    bullets=["Role: UX Designer", "Timeline: 7 months; Launched December 2025", "Team: Kindle UX, PM, and engineering teams"],
    notes="This was Kindle's first color writing device ever. I owned the entire color writing experience - annotations, writing tools, bookmarks, and templates.")

add_content_slide("Opportunity", "Kindle's first color writing device",
    text="Kindle Color Scribe introduces color as a new medium for annotation, writing, and organization.",
    bullets=["Define what color means on a device built around focus and simplicity", "Enable richer ways to highlight, write, and structure ideas", "Preserve the calm, immersive reading experience Kindle is known for"],
    notes="The opportunity was exciting but also challenging - how do you add color without making the device feel cluttered or distracting?")

add_content_slide("Design Challenge", "How might we introduce powerful color writing tools while still preserving the flow of deep reading and writing?",
    notes="This tension between power and simplicity was the core challenge throughout the project.")

add_content_slide("Design Challenge #1", "Adding color to annotations",
    text="The in-book contextual action bar (CAB) had to feel almost invisible—powerful when needed, effortless when not.",
    bullets=["Challenge: Fitting more functionality into constrained, frequently used space", "Goal: Minimize cognitive load while preserving fast, one-tap access"],
    images="assets/Project 2/image 8.png",
    notes="The CAB is one of the most-used UI elements on Kindle. Any changes here had to be extremely well-considered.")

add_image_slide("", "Iterating on the CAB",
    "assets/Project 2/image 9.png",
    notes="Early explorations tested multiple interaction models - toggles, nested menus, explicit delete actions. We used design crits and usability principles to pressure-test each option.")

add_content_slide("", "Simplifying the menu",
    bullets=["Removed explicit delete action—folded into 'None' option in color submenu", "Kept highlight and underline as first-class actions (not nested)", "Prioritized clarity for today's core use cases"],
    images="assets/Project 2/image 10.png",
    notes="One of the most important decisions was removing the explicit delete action. This reframed deletion as part of editing, not a separate mode.")

add_content_slide("", "Designing cross-surface",
    text="The CAB needed to flex across color and black-and-white devices with different technical constraints.",
    bullets=["Surface-appropriate designs preserving same mental model", "Adapting interaction patterns per device capabilities"],
    images="assets/Project 2/image 11.png",
    notes="Rather than enforcing strict uniformity, I advocated for surface-appropriate designs. This added engineering complexity but resulted in better UX for each user group.")

add_content_slide("Result", "Final annotation experience",
    text="Fast, flexible annotation without interrupting reading flow.",
    bullets=["CSAT: 6.14/7 for color highlights and underlines in beta testing", "Annotations not flagged as pain point during beta or launch"],
    images="assets/Project 2/CAB.gif",
    notes="The high CSAT score validated our approach. Most importantly, users could annotate without breaking immersion.")

add_content_slide("Design Challenge #2", "Color switching friction",
    text="User research and beta feedback consistently surfaced color switching as a pain point.",
    bullets=["Color options nested in submenu requiring multiple taps", "Most users only use 2-3 colors in a session"],
    images="assets/Project 2/image 13.png",
    notes="This was identified through beta feedback. The existing design worked functionally but interrupted flow for users who switch colors frequently.")

add_content_slide("", "Exploring solutions",
    bullets=["Color slots: Fastest switching, minimal cognitive load once learned", "Pen presets: Bundled tool type + color, but added setup complexity", "Floating palettes: Most discoverable, but covered writing space"],
    images="assets/Project 2/image 14.png",
    notes="Each direction had clear tradeoffs. We evaluated based on speed, simplicity, and technical feasibility given our timeline.")

add_content_slide("Result", "Color slots solution",
    text="Aligned on color slots as the best balance of speed, simplicity, and feasibility.",
    bullets=["Customer satisfaction: 6.10/7", "Over 80% rated pen color controls as easy or very easy to use", "Writing in color felt as fluid as black ink"],
    images="assets/Project 2/color switching gif.gif",
    notes="Internal beta validated this approach. The key success metric was that color didn't slow users down.")

add_content_slide("Design Challenge #3", "Templates",
    text="Rather than designing from taste, I partnered with UX Research to anchor decisions in user input.",
    images="assets/Project 2/image 18.png",
    notes="Templates are inherently personal, so we needed research to guide decisions rather than relying on internal opinions.")

add_content_slide("", "What users actually wanted",
    bullets=["Users favored simple, versatile templates over decorative designs", "Contrary to internal pressure, users preferred black-and-white templates", "Customizability emerged as strongest unmet need"],
    images=["assets/Project 2/image 19.png", "assets/Project 2/image 20.png"],
    notes="This was surprising - there was internal pressure to add more colorful templates, but users consistently preferred simplicity.")

add_content_slide("", "Making 30 templates manageable",
    text="As the library grew, discoverability became the next challenge.",
    bullets=["Solution: 'Recently used' row at top of picker", "93% rating for creating notebooks with templates", "Users reported it as 'intuitive'"],
    images="assets/Project 2/image 25.png",
    notes="The recently used row optimized for real usage patterns without adding UI complexity.")

add_image_slide("Solution", "Writing that feels like pen and paper",
    "assets/Project 2/image 1.png",
    notes="Configurable color slots enable one-tap switching between most-used colors, while lasso select helps make quick bulk color adjustments.")

add_image_slide("", "Powerful annotation experience",
    "assets/Project 2/image 2.png",
    notes="Highlight, underline, and annotate in color. Color and date filters let you quickly find what matters most.")

add_image_slide("", "Curated templates",
    ["assets/Project 2/image 21.png", "assets/Project 2/image 22.png", "assets/Project 2/image 23.png"],
    notes="Added 12 new, top-ranked templates based on user research, plus familiar colored options like college ruled and legal pad.")

add_content_slide("Outcome", "Launch",
    text="Kindle Scribe Colorsoft launched December 10, 2025.",
    bullets=['"Smooth writing feel. Easy to use." — Wired', '"The software stays out of the way" — Inc.', '"A feature I didn\'t know I wanted or needed." — Good Housekeeping'],
    images="assets/Project 2/image 27.png",
    notes="The press reception was very positive. The key theme in reviews was that the software felt natural and unobtrusive.")

add_content_slide("Learnings", "What I learned",
    bullets=["Define new interaction patterns while earning user trust on traditionally grayscale product", "Systems thinking across surfaces—when to unify vs. allow divergence", "Reduce cognitive load and preserve flow in high-speed interactions", "Treat technical constraints as design inputs—adapting solutions to remain intuitive"],
    notes="This project pushed me to think about systems design and how to introduce new capabilities without overwhelming users. The biggest lesson was treating constraints as design inputs rather than blockers.")

# CLOSING
add_title_slide("Thank you!", "Elise Nguyen\nelisenguyen.com",
    notes="Thank them for their time. Open for questions.")

# Save
prs.save('portfolio_slides.pptx')
print("Created portfolio_slides.pptx with images, styling, and speaker notes!")
